# coding=utf-8

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from get_data import gxj_chart


def analyze_ppt(input_file, output_file):
    """
    通过模板文件(input_file)生成一个ouput_file,
    以方便查看模板中placeholder的编号。
    """
    prs = Presentation(input_file)

    for i, layout in enumerate(prs.slide_layouts):
        # 用当前板式(layout)新建一张幻灯片(slide)
        slide = prs.slides.add_slide(layout)

        # 是否有标题占位符
        try:
            title = slide.shapes.title
            title.text = f'page-{i}:title'
        except AttributeError:
            print(f'>>> page-{i}:no title')

        # 将其他占位符(placeholders)命名为layout-placeholder编号
        for each in slide.placeholders:
            each.text = f'layouts[{i}-placeholder[{each.placeholder_format.idx}]]'

    prs.save(output_file)


def insert_text(slide, index, text):
    """在传入的slide中找到指定编号的placeholder,并插入文字
    :param slide: 幻灯片对象
    :param index: placeholder在当前slide上的编号
    :param text: 插入文字
    """
    slide.placeholders[index].text = text


def insert_picture(slide, index, image):
    """
    :param slide: 幻灯片对象
    :param index: placeholder在当前slide上的编号
    :param image_path: 图片地址
    """
    slide.placeholders[index].insert_picture(image)


def insert_gxj_chart(slide, index, df, color, chart_file, x_style={}):
    """
    在指定placeholder 中插入 gxj_chart
    :param slide:  slide对象
    :param index: placeholder编号
    :param df: DataFrame对象
    :param color: 图板本色方案，如('#494949', '#C10F2A', '#F39617')
    :param chart_file: 图片文件路径，其实设什么都无所谓
    """
    print(f'>>> 正在插入图片...')
    # 生成chart文件
    gxj_chart(df, color=color, output=chart_file, xticks_style=x_style)
    # 将文件插入placeholder
    insert_picture(slide, index, chart_file)


def df_to_table(slide, index, df):
    """将df插入到指定的表格占位符中
    :param slide: slide对象
    :param index: placeholder编号
    :param df: DataFrame对象
    """
    placeholder = slide.placeholders[index]
    # 确定表格行、列数
    rows, cols = df.shape
    # 插入表格，多一行供填写表头
    tb = placeholder.insert_table(rows + 1, cols).table

    # 调整列宽
    columns = tb.columns
    columns[0].width = Inches(0.7)
    columns[1].width = Inches(2.5)
    columns[2].wdith = Inches(1)
    columns[3].wdith = Inches(1.5)

    # 填写表头
    colnames = list(df.columns)
    for i, col_name in enumerate(colnames):
        cell = tb.cell(0, i)
        cell.text = col_name
        # 垂直居中，具体见pptx.enum.text.MSO_VERTICAL_ANCHOR源代码
        cell.vertical_anchor = 3
        # 每个单元格应该只有1段，改写该段的字号与对齐方式
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.size = Pt(12)
        paragraph.alignment = PP_ALIGN.CENTER

    # 填写数据
    m = df.as_matrix()
    for row in range(rows):
        for col in range(cols):
            value = m[row, col]
            cell = tb.cell(row + 1, col)
            cell.text = str(value)
            # 样式
            cell.vertical_anchor = 3
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(12)
            paragraph.alignment = PP_ALIGN.CENTER