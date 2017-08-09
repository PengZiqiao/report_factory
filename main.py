# coding=utf-8
from pptx import Presentation
from winsun.date import Week
from winsun.tools import BIE_SHU

from get_data import PyGis, ShuoLi
from ppt import insert_text, insert_gxj_chart, df_to_table


def liangjia(wuye):
    """
    创建一页市场量价的幻灯片
    :param prs: Presentation对象
    :param wuye:
    :return:
    """
    # 常量
    df = dict()
    chart_color = ('#4A4949', '#C10F2A', '#F39617')
    chart_file = base_path + 'chart.png'

    # 根据模板中 样式0 创建1页幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    # 查询, 数据来源若为excel， 须改写以下代码
    print(f'>>> 正在查询{wuye}走势...')
    df[f'{wuye}走势'] = g.trend_gxj(usage[wuye])
    print(f'>>> 正在查询{wuye}分板块量价...')
    df[f'{wuye}板块'] = g.plate_gxj(usage[wuye])

    # 标题
    insert_text(slide, 0, f'{wuye}市场-量价')
    # 图表标题
    insert_text(slide, 11, f'南京（不含高淳溧水）近10周{wuye}市场供销量价')
    insert_text(slide, 13, f'2017年第{wN}周南京（不含高淳溧水）{wuye}市场分板块供销量价')
    # 插入图表 x_style来控制横坐标轴样式
    insert_gxj_chart(slide, 12, df[f'{wuye}走势'], chart_color, chart_file, x_style={'fontsize': 10})
    insert_gxj_chart(slide, 14, df[f'{wuye}板块'], chart_color, chart_file)
    # 说理
    shuoli = ShuoLi(df[f'{wuye}走势']).all()
    # 商办说理中添加商办销量前三
    if wuye == '商业' or wuye == '办公':
        else_info = f'本周成交面积榜前三：{g.rank_shuoli(wuye)}'
    else:
        else_info = '补充说理。'
    insert_text(slide, 15, f'本周{wuye}市场{shuoli}\r{else_info}')


def zhuzhai_rank():
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    # 标题
    insert_text(slide, 0, '住宅市场-排行')
    # 日期
    year = w.sunday.year
    monday = w.monday
    monday = f'{monday.month}月{monday.day}日'
    sunday = w.sunday
    sunday = f'{sunday.month}月{sunday.day}日'
    date_text = f'{year}年第{wN}周（{monday}-{sunday}）'
    # 表格标题
    insert_text(slide, 11, f'{date_text}成交面积排行榜')
    insert_text(slide, 13, f'{date_text}成交套数排行榜')
    # 插入表格
    print('>>> 正在查询排行榜数据...')
    df_mj, df_ts = g.zhuzhai_rank()
    df_to_table(slide, 14, df_mj)
    df_to_table(slide, 15, df_ts)


if __name__ == '__main__':
    # 路径
    base_path = 'e:/周报测试/'

    # 日期
    w = Week()
    wN = w.N

    # gis系统改版后，PyGis与GisSpider须重写
    g = PyGis()

    # 通过模板创建ppt
    usage = {'住宅': ['住宅'], '商业': ['商业'], '办公': ['办公'], '别墅': BIE_SHU}
    prs = Presentation('template.pptx')

    # 遍历4种物业类型
    for each in usage:
        print(f'====={each}市场=====')
        liangjia(each)
        # 住宅排行
        if each == '住宅':
            print(f'=====住宅排行=====')
            zhuzhai_rank()

    # 保存ppt
    print('>>> 正在保存ppt...')
    prs.save(base_path + f'2017年第{wN}周周报.pptx')
    print('>>> 保存完成！')
